#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成智能体搭建PPT（2页）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色
DARK_BLUE = RGBColor(0x0f, 0x17, 0x3d)
ACCENT_BLUE = RGBColor(0x20, 0x5e, 0xc8)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x5c, 0x5c, 0x5c)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
ACCENT = RGBColor(0x00, 0xc9, 0x5f)
RED = RGBColor(0xe5, 0x3e, 0x3e)
PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
ORANGE = RGBColor(0xff, 0x98, 0x00)
CYAN = RGBColor(0x06, 0xb6, 0xd4)


def new_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.85))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12), Inches(0.55))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return slide


def add_box(x, y, w, h, text, color, text_color=None, size=11):
    slide = prs.slides[-1]
    if text_color is None:
        text_color = WHITE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    t = slide.shapes.add_textbox(Inches(x), Inches(y + (h - 0.3)/2), Inches(w), Inches(0.35))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return box


def add_card(x, y, w, h, icon, title, items, color):
    slide = prs.slides[-1]
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = WHITE
    c.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    i = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(w - 0.3), Inches(0.35))
    tf = i.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    cont = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.5), Inches(w - 0.3), Inches(h - 0.6))
    tf = cont.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸ " + item
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY


def add_arrow(x, y, text="▶", size=18):
    slide = prs.slides[-1]
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.5), Inches(0.4))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER


# ========== 第1页：什么是AI智能体 ==========
new_slide("AI智能体：让AI真正行动起来")

# 核心概念
add_card(0.3, 1.05, 12.7, 1.5, "🤖", "AI智能体（Agent）",
    ["传统AI：只能对话，你问它答，被动响应", "AI智能体：自主决策，调用工具，执行任务", "核心能力：感知→决策→行动→反馈，持续循环直到完成任务"], ACCENT_BLUE)

# 智能体 vs 对话
add_card(0.3, 2.7, 6.2, 2.3, "💬", "传统AI对话",
    ["用户提问", "AI回答", "结束", "被动响应，等待指令"], RED)

add_card(6.8, 2.7, 6.2, 2.3, "🛡️", "AI智能体",
    ["用户下达目标", "AI自主规划", "调用工具执行", "持续反馈优化直到完成"], ACCENT)

# 底部流程
t = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(4), Inches(0.3))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "🔄 智能体工作循环"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# 四个循环步骤
steps = [
    (0.5, "🎯", "感知", "理解任务\n分析目标"),
    (3.4, "🧠", "决策", "规划步骤\n选择工具"),
    (6.3, "⚙️", "执行", "调用工具\n获取结果"),
    (9.2, "📊", "反馈", "分析结果\n优化迭代"),
]

for x, icon, title, desc in steps:
    add_box(x, 5.6, 2.7, 1.5, f"{icon}\n{title}\n{desc}", ACCENT_BLUE, WHITE, 10)
    if x < 9.2:
        add_arrow(x + 2.75, 6.2, "⟳", 22)


# ========== 第2页：智能体工具与执行 ==========
new_slide("智能体执行流程：工具调用实战")

# 顶部说明
t = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.35))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "用户上传演练文档 → AI分析并制定计划 → Agent循环执行工具 → 生成报告"
p.font.size = Pt(12)
p.font.color.rgb = GRAY

# 左侧：可用工具
add_card(0.3, 1.45, 4.5, 5.7, "🛠️", "可用工具集",
    [
        "🌐 network_scan - 网络扫描",
        "   发现存活主机、端口、服务",
        "📸 web_screenshot - Web截图",
        "   可视化Web界面",
        "🔓 bruteforce_ssh - SSH检测",
        "   SSH弱口令检测",
        "🔓 bruteforce_rdp - RDP检测",
        "   RDP弱口令检测",
        "🔓 bruteforce_mysql - MySQL检测",
        "   数据库弱口令检测",
        "🐝 honeypot_audit - 蜜罐审计",
        "   查询蜜罐攻击日志",
        "📋 generate_report - 生成报告",
        "   汇总演练结果",
    ], ACCENT_BLUE)

# 右侧：执行流程
t2 = prs.slides[-1].shapes.add_textbox(Inches(5.1), Inches(1.45), Inches(4), Inches(0.35))
tf = t2.text_frame
p = tf.paragraphs[0]
p.text = "🔄 Agent执行流程"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# 流程图
flow_steps = [
    (5.1, "📄", "上传文档", "用户上传\n演练文档", ACCENT),
    (5.1, 2.5, "🧠", "AI分析", "提取目标\n制定计划", PURPLE),
    (5.1, 4.0, "⚡", "执行工具", "调用network_scan\n扫描网络", ORANGE),
    (5.1, 5.5, "📸", "截图取证", "web_screenshot\n记录界面", CYAN),
    (5.1, 7.0, "🔓", "弱口检测", "bruteforce_*\n检测弱口令", RED),
]

# 用方块表示流程
y_pos = 1.9
flow_blocks = [
    ("📄", "上传文档", "用户上传演练文档", ACCENT),
    ("🧠", "AI分析", "提取目标网络/服务", PURPLE),
    ("⚡", "执行network_scan", "发现存活主机", ORANGE),
    ("📸", "执行web_screenshot", "Web界面截图", CYAN),
    ("🔓", "执行bruteforce_*", "检测SSH/RDP/MySQL", RED),
    ("🐝", "执行honeypot_audit", "查询蜜罐日志", ACCENT_BLUE),
    ("📋", "生成报告", "汇总演练结果", ACCENT),
]

for i, (icon, title, desc, color) in enumerate(flow_blocks):
    y = 1.85 + i * 0.78
    # 编号圆
    circ = prs.slides[-1].shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.1), Inches(y), Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()

    n = prs.slides[-1].shapes.add_textbox(Inches(5.1), Inches(y + 0.05), Inches(0.4), Inches(0.35))
    tf = n.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 标题
    ti = prs.slides[-1].shapes.add_textbox(Inches(5.6), Inches(y), Inches(3.2), Inches(0.25))
    tf = ti.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # 描述
    de = prs.slides[-1].shapes.add_textbox(Inches(5.6), Inches(y + 0.28), Inches(3.2), Inches(0.3))
    tf = de.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY

    # 箭头
    if i < len(flow_blocks) - 1:
        add_arrow(5.15, y + 0.6, "↓", 14)

# 右侧：案例展示
add_card(9.2, 1.45, 3.8, 5.7, "🎯", "实际案例",
    [
        "演练文档：",
        "目标：192.168.1.0/24",
        "重点：Web、SSH服务",
        "",
        "AI自动执行：",
        "1. 网络扫描→发现3台主机",
        "2. Web截图→2个Web服务",
        "3. SSH检测→1个弱口令",
        "4. 蜜罐审计→2条记录",
        "5. 生成报告→完成",
        "",
        "整个过程：AI自主完成",
        "人工介入：仅上传文档",
    ], ORANGE)


# 保存
prs.save("C:/Users/lzx78/Desktop/AimiGuard/ai_agent_intro.pptx")
print("PPT已生成!")
