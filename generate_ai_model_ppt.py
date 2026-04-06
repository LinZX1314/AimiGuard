#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成AI模型介绍PPT（2页）- 流程图展示
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色
DARK_BLUE = RGBColor(0x0f, 0x17, 0x3d)
ACCENT_BLUE = RGBColor(0x20, 0x5e, 0xc8)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x5c, 0x5c, 0x5c)
LIGHT_GRAY = RGBColor(0xf5, 0xf5, 0xf5)
ACCENT = RGBColor(0x00, 0xc9, 0x5f)
RED = RGBColor(0xe5, 0x3e, 0x3e)
PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
ORANGE = RGBColor(0xff, 0x98, 0x00)


def new_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.85))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12), Inches(0.55))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return slide


def add_box(x, y, w, h, text, color, text_color=None, size=11):
    slide = prs.slides[-1]
    if text_color is None:
        text_color = WHITE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    t = slide.shapes.add_textbox(Inches(x), Inches(y + (h - 0.3)/2), Inches(w), Inches(0.35))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return box


def add_arrow(x, y, text="▶"):
    slide = prs.slides[-1]
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.5), Inches(0.4))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER


# ========== 第1页：训练流程 ==========
new_slide("AI安全专家训练流程")

# 顶部说明
t = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.4))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "基础模型：通义千问 Qwen3.5-9B（通用大模型）  →  经过专精化训练  →  网络安全AI专家"
p.font.size = Pt(13)
p.font.color.rgb = GRAY

# ===== 第一阶段：数据准备 =====
t1 = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2), Inches(0.35))
tf = t1.text_frame
p = tf.paragraphs[0]
p.text = "📊 数据准备"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# 三个数据源
y1 = 1.95
add_box(0.5, y1, 2.5, 1.0, "🔐\n漏洞库\nCVE/CNVD", ACCENT_BLUE, WHITE, 10)
add_arrow(3.1, y1 + 0.3)
add_box(3.5, y1, 2.5, 1.0, "📊\n实战日志\nNmap/蜜罐记录", ACCENT, WHITE, 10)
add_arrow(6.1, y1 + 0.3)
add_box(6.5, y1, 2.5, 1.0, "📚\n安全文档\nRFC/OWASP", PURPLE, WHITE, 10)

# 合并箭头
add_arrow(9.15, y1 + 0.3)
add_box(9.6, y1, 3.3, 1.0, "✅\n高质量训练数据", ORANGE, WHITE, 11)

# ===== 第二阶段：LoRA微调 =====
t2 = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(2), Inches(0.35))
tf = t2.text_frame
p = tf.paragraphs[0]
p.text = "⚙️ LoRA微调"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

y2 = 3.65
add_box(0.5, y2, 3.5, 1.3, "🧠\nLoRA技术\n只更新0.1%参数\n保留原有能力", ACCENT_BLUE, WHITE, 10)

add_arrow(4.15, y2 + 0.45)

add_box(4.6, y2, 3.0, 1.3, "💰\n低成本训练\n一块RTX 4090\n即可完成", ACCENT, WHITE, 10)

add_arrow(7.75, y2 + 0.45)

add_box(8.2, y2, 4.6, 1.3, "🛡️\n专家能力\n威胁解读 / 漏洞分析\n应急响应", PURPLE, WHITE, 10)

# ===== 第三阶段：量化压缩 =====
t3 = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(2), Inches(0.35))
tf = t3.text_frame
p = tf.paragraphs[0]
p.text = "📦 量化压缩"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

y3 = 5.55

# FP16
add_box(0.5, y3, 2.2, 1.1, "FP16\n18GB显存\n服务器级", RED, WHITE, 10)

add_arrow(2.85, y3 + 0.35, "→")

# GPTQ
add_box(3.2, y3, 2.2, 1.1, "GPTQ\n量化压缩\nFP16→INT4", ACCENT_BLUE, WHITE, 10)

add_arrow(5.55, y3 + 0.35, "→")

# INT4
add_box(5.9, y3, 2.2, 1.1, "INT4\n4~6GB显存\n游戏显卡跑", ACCENT, WHITE, 10)

add_arrow(8.25, y3 + 0.35, "→")

# 最终模型
add_box(8.7, y3, 4.1, 1.1, "🚀\n网络安全AI专家\n轻量化 / 高效率 / 低成本", ORANGE, WHITE, 11)

# 底部
bottom = prs.slides[-1].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5))
bottom.fill.solid()
bottom.fill.fore_color.rgb = DARK_BLUE
bottom.line.fill.background()

t = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.4))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "部署至 AimiGuard AI对话模块 →  成为网络安全AI助手  →  赋能中小企业安全运营"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER


# ========== 第2页：能力展示 ==========
new_slide("AI安全专家能力展示")

# 三大核心能力
t = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(4), Inches(0.35))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "🎯 核心能力"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# 三个能力卡片
cards = [
    (0.5, 1.45, "🚨", "威胁解读", [
        "输入告警日志",
        "识别攻击类型",
        "判断处置方式",
        "输出分析报告"
    ], RED),
    (4.55, 1.45, "🔧", "漏洞分析", [
        "输入CVE编号",
        "解读漏洞原理",
        "分析利用条件",
        "给出修复建议"
    ], ACCENT),
    (8.6, 1.45, "⚡", "应急响应", [
        "遇到安全事件",
        "给出排查思路",
        "提供处置步骤",
        "辅助快速响应"
    ], ORANGE),
]

for x, y, icon, title, items, color in cards:
    # 卡片背景
    c = prs.slides[-1].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.85), Inches(2.4))
    c.fill.solid()
    c.fill.fore_color.rgb = WHITE
    c.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)

    # 顶部色条
    bar = prs.slides[-1].shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.85), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # 图标
    ic = prs.slides[-1].shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.2), Inches(3.5), Inches(0.5))
    tf = ic.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # 内容
    cont = prs.slides[-1].shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.7), Inches(3.5), Inches(1.6))
    tf = cont.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸ " + item
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

# 与AI对话集成
t2 = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(6), Inches(0.35))
tf = t2.text_frame
p = tf.paragraphs[0]
p.text = "🔗 与AI对话模块集成"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# 集成示意
integ_bg = prs.slides[-1].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.2))
integ_bg.fill.solid()
integ_bg.fill.fore_color.rgb = WHITE
integ_bg.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)

# 用户
add_box(0.8, 4.65, 2.0, 0.9, "👤\n用户提问", ACCENT_BLUE, WHITE, 10)

add_arrow(2.95, 4.85)

# AI对话
add_box(3.4, 4.65, 2.5, 0.9, "💬\nAI对话模块\n智能交互", ACCENT, WHITE, 10)

add_arrow(6.05, 4.85)

# 安全专家
add_box(6.5, 4.65, 2.5, 0.9, "🛡️\n安全专家模型\n威胁分析", PURPLE, WHITE, 10)

add_arrow(9.15, 4.85)

# 结果
add_box(9.6, 4.65, 2.8, 0.9, "📋\n专业建议\n自动封禁/处置", ORANGE, WHITE, 10)

# 四大优势
t3 = prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(4), Inches(0.35))
tf = t3.text_frame
p = tf.paragraphs[0]
p.text = "✨ 核心优势"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

advantages = [
    (0.5, "🛡️", "不知疲倦", "7×24在线", ACCENT_BLUE),
    (3.45, "🎯", "精准专业", "领域专家", ACCENT),
    (6.4, "⚡", "快速响应", "秒级分析", PURPLE),
    (9.35, "💰", "成本低廉", "显卡即可", ORANGE),
]

for x, icon, title, desc, color in advantages:
    add_box(x, 6.3, 2.85, 1.0, f"{icon} {title}\n{desc}", color, WHITE, 10)


# 保存
prs.save("C:/Users/lzx78/Desktop/AimiGuard/ai_model_2pages.pptx")
print("PPT已生成!")
